from pptx import Presentation

prs = Presentation("C:\Thameem\Testing\xxx\Account Onboarding Deck V2.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
print("1111",text_runs)